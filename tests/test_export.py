from pathlib import Path

import nbformat as nbf
import pytest
from pptx import Presentation

from models.schemas import (
    Assignment,
    CodeExercise,
    CodeExerciseSet,
    CourseOutline,
    DiagramSet,
    DiagramSpec,
    Lesson,
    NotebookCell,
    NotebookDraft,
    PersonaProfile,
    Quiz,
    QuizQuestion,
    SpeakerNoteSection,
    SpeakerNotes,
)
from tools.export import export_run
from tools.notebook_export import build_notebook, export_notebook_to_file
from tools.pptx_export import LessonExtras, build_pptx, export_pptx_to_file


@pytest.fixture
def sample_draft():
    return NotebookDraft(
        lesson_title="Ownership",
        cells=[
            NotebookCell(cell_type="markdown", content="# Ownership"),
            NotebookCell(cell_type="code", content="x = 1"),
        ],
    )


@pytest.fixture
def sample_course():
    return CourseOutline(
        course_title="Rust Basics",
        lessons=[
            Lesson(title="Ownership", objectives=["explain ownership"], estimated_minutes=20),
            Lesson(title="Borrowing", objectives=["explain borrowing"], estimated_minutes=25),
        ],
    )


# --- Notebook export ----------------------------------------------------------


def test_build_notebook_minimal_is_valid(sample_draft):
    nb = build_notebook(sample_draft)
    nbf.validate(nb)  # raises on malformed notebook
    assert len(nb["cells"]) == 2


def test_build_notebook_includes_all_extras(sample_draft):
    exercises = CodeExerciseSet(
        lesson_title="Ownership",
        exercises=[CodeExercise(prompt="Do X", starter_code="pass", solution_code="return 1")],
    )
    diagrams = DiagramSet(
        lesson_title="Ownership",
        diagrams=[DiagramSpec(title="Flow", mermaid_code="graph TD; A-->B;", description="shows flow")],
    )
    quiz = Quiz(
        lesson_title="Ownership",
        questions=[
            QuizQuestion(question="What?", options=["A", "B", "C", "D"], correct_answer_index=2, explanation="C is right")
        ],
    )
    assignment = Assignment(
        lesson_title="Ownership",
        title="Build something",
        instructions="Do the thing",
        deliverables=["a file"],
        rubric=["works"],
    )

    nb = build_notebook(sample_draft, code_exercises=exercises, diagrams=diagrams, quiz=quiz, assignment=assignment)
    nbf.validate(nb)

    full_text = "\n".join(c["source"] for c in nb["cells"])
    assert "Flow" in full_text
    assert "Do X" in full_text
    assert "What?" in full_text
    assert "Build something" in full_text
    # Answer key should be present but not the raw index leaking unexplained
    assert "Answer key" in full_text


def test_export_notebook_to_file_writes_readable_file(tmp_path, sample_draft):
    nb = build_notebook(sample_draft)
    out = export_notebook_to_file(nb, tmp_path / "sub" / "course.ipynb")
    assert out.exists()
    reloaded = nbf.read(out, as_version=4)
    assert len(reloaded["cells"]) == 2


# --- PPTX export ---------------------------------------------------------------


def test_build_pptx_basic_structure(sample_course):
    prs = build_pptx(sample_course)
    # title + agenda + one slide per lesson (2) + thanks = 5
    assert len(prs.slides) == 5
    assert prs.slides[0].shapes.title.text == "Rust Basics"


def test_build_pptx_includes_persona_subtitle(sample_course):
    persona = PersonaProfile(audience_level="beginner", tone="friendly", prerequisites=[], learning_goals=["x"])
    prs = build_pptx(sample_course, persona=persona)
    subtitle = prs.slides[0].placeholders[1].text
    assert "beginner" in subtitle


def test_build_pptx_sets_speaker_notes(sample_course):
    notes = SpeakerNotes(lesson_title="Ownership", sections=[SpeakerNoteSection(section_title="Intro", script="Hello!")])
    prs = build_pptx(sample_course, lesson_extras={0: LessonExtras(speaker_notes=notes)})
    # slide 0 = title, slide 1 = agenda, slide 2 = lesson 0 ("Ownership")
    lesson_slide = prs.slides[2]
    assert "Hello!" in lesson_slide.notes_slide.notes_text_frame.text


def test_build_pptx_adds_quiz_slides(sample_course):
    quiz = Quiz(
        lesson_title="Ownership",
        questions=[QuizQuestion(question="Q1?", options=["A", "B", "C", "D"], correct_answer_index=0, explanation="A")],
    )
    prs_without = build_pptx(sample_course)
    prs_with = build_pptx(sample_course, lesson_extras={0: LessonExtras(quiz=quiz)})
    assert len(prs_with.slides) == len(prs_without.slides) + 1


def test_export_pptx_to_file_writes_openable_file(tmp_path, sample_course):
    prs = build_pptx(sample_course)
    out = export_pptx_to_file(prs, tmp_path / "sub" / "slides.pptx")
    assert out.exists()
    reopened = Presentation(str(out))
    assert len(reopened.slides) == len(prs.slides)


# --- Full run export ------------------------------------------------------------


def test_export_run_requires_curriculum(tmp_path, monkeypatch):
    monkeypatch.setattr("tools.export.settings.outputs_dir", tmp_path)
    with pytest.raises(ValueError, match="curriculum"):
        export_run("fake-run-id", {"topic": "x"})


def test_export_run_produces_zip_with_notebook_and_pptx(tmp_path, monkeypatch, sample_course, sample_draft):
    monkeypatch.setattr("tools.export.settings.outputs_dir", tmp_path)
    context = {"curriculum": sample_course, "notebook": sample_draft}

    zip_path = export_run("test-run-123", context)

    assert zip_path.exists()
    assert zip_path.suffix == ".zip"

    import zipfile

    with zipfile.ZipFile(zip_path) as zf:
        names = zf.namelist()
    assert "course.ipynb" in names
    assert "course_slides.pptx" in names


def test_export_run_skips_notebook_without_notebook_output(tmp_path, monkeypatch, sample_course):
    monkeypatch.setattr("tools.export.settings.outputs_dir", tmp_path)
    context = {"curriculum": sample_course}  # no "notebook" key

    zip_path = export_run("test-run-456", context)

    import zipfile

    with zipfile.ZipFile(zip_path) as zf:
        names = zf.namelist()
    assert "course.ipynb" not in names
    assert "course_slides.pptx" in names


@pytest.mark.parametrize(
    "bad_run_id",
    ["../escape", "..\\escape", "foo/bar", "foo\\bar", "", "has spaces", "a" * 200],
)
def test_export_run_rejects_unsafe_run_id(tmp_path, monkeypatch, sample_course, bad_run_id):
    monkeypatch.setattr("tools.export.settings.outputs_dir", tmp_path)
    with pytest.raises(ValueError, match="run_id"):
        export_run(bad_run_id, {"curriculum": sample_course})
