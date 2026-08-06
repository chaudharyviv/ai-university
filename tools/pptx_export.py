"""
PPTX export - turns the course outline (+ optional per-lesson extras)
into a slide deck via python-pptx.

Since Phase 3/4 agents currently only generate content for lesson 0,
`lesson_extras` covers that one lesson richly (speaker notes as real
slide notes, quiz recap, diagram descriptions); every other lesson gets
a plain title + objectives slide. Looping full extras across every
lesson is a natural Phase 6+ extension once agents run per-lesson.
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from loguru import logger
from pptx import Presentation

from models.schemas import CourseOutline, DiagramSet, PersonaProfile, Quiz, SpeakerNotes


@dataclass
class LessonExtras:
    """Extra content available for one lesson (by index) beyond title/objectives."""

    diagrams: DiagramSet | None = None
    quiz: Quiz | None = None
    speaker_notes: SpeakerNotes | None = None


def _add_title_slide(prs: Presentation, title: str, subtitle: str = "") -> None:
    layout = prs.slide_layouts[0]  # "Title Slide"
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle


def _add_bullet_slide(prs: Presentation, title: str, bullets: list[str]):
    layout = prs.slide_layouts[1]  # "Title and Content"
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, bullet in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = bullet
        p.level = 0
    return slide


def _set_notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


def build_pptx(
    course: CourseOutline,
    *,
    persona: PersonaProfile | None = None,
    lesson_extras: dict[int, LessonExtras] | None = None,
) -> Presentation:
    lesson_extras = lesson_extras or {}
    prs = Presentation()

    subtitle = f"Audience: {persona.audience_level}" if persona else ""
    _add_title_slide(prs, course.course_title, subtitle)

    # Agenda slide
    _add_bullet_slide(prs, "Agenda", [lesson.title for lesson in course.lessons])

    for i, lesson in enumerate(course.lessons):
        extras = lesson_extras.get(i, LessonExtras())

        slide = _add_bullet_slide(
            prs,
            lesson.title,
            [f"{obj}" for obj in lesson.objectives] + [f"~{lesson.estimated_minutes} min"],
        )
        if extras.speaker_notes and extras.speaker_notes.sections:
            script = "\n\n".join(
                f"[{s.section_title}]\n{s.script}" for s in extras.speaker_notes.sections
            )
            _set_notes(slide, script)

        if extras.diagrams and extras.diagrams.diagrams:
            diagram_bullets = [f"{d.title}: {d.description}" for d in extras.diagrams.diagrams]
            _add_bullet_slide(prs, f"{lesson.title} - Diagrams", diagram_bullets)
            # Note: diagrams are included as text descriptions, not rendered
            # images - rendering Mermaid to a real graphic would need an
            # external renderer (e.g. mermaid-cli), out of scope for now.

        if extras.quiz and extras.quiz.questions:
            for q in extras.quiz.questions:
                _add_bullet_slide(prs, "Quick Check", [q.question] + [f"  {o}" for o in q.options])

    _add_bullet_slide(prs, "Thanks!", ["Questions?"])
    return prs


def export_pptx_to_file(prs: Presentation, path: Path) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    logger.info("Wrote slide deck: {} ({} slides)", path, len(prs.slides))
    return path
